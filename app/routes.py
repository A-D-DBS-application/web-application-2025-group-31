from flask import Blueprint, render_template, request, redirect, url_for, session, Response, jsonify
from werkzeug.security import generate_password_hash, check_password_hash
from app import db
from app.models import AppUser, Company, Metric, AuditLog, ChangeEvent, MetricHistory, Sector
from decimal import Decimal
import csv
import io
from app.scraper import scrape_website
from datetime import datetime
from app.similarity import top_similar_companies
from app.auth import login_required
from app.auth import admin_required

bp = Blueprint('main', __name__)

METRIC_OPTIONS = ["Pricing", "Features", "Reviews", "Funding", "Hiring"]

# ======================================================
# TEKST NORMALISATIE EN VERGELIJKING
# ======================================================

import difflib
import re

def normalize_text(s: str):
    if not s:
        return ""
    s = s.lower().strip()
    s = re.sub(r"\s+", " ", s)
    return s

def texts_similar(a: str, b: str, threshold=0.90):
    """Return True if texts are almost the same (ignore small AI noise)."""
    a_norm = normalize_text(a)
    b_norm = normalize_text(b)
    if not a_norm and not b_norm:
        return True
    similarity = difflib.SequenceMatcher(None, a_norm, b_norm).ratio()
    return similarity >= threshold

def normalize_list(lst):
    if not lst:
        return []
    return sorted([normalize_text(x) for x in lst if x and str(x).strip() != ""])


# =====================================================
# HELPER FUNCTIES (belangrijk voor correcte Supabase types)
# =====================================================

def safe_int(x):
    """Zet string of AI-output om naar integer of None."""
    try:
        return int(str(x).replace(",", "").strip())
    except:
        return None


def safe_float(x):
    """Zet currency/tekst om naar float (Supabase NUMERIC)."""
    if x is None:
        return None
    try:
        cleaned = (
            str(x)
            .replace("€", "")
            .replace("$", "")
            .replace(",", "")
            .replace("m", "000000")
            .replace("M", "000000")
            .strip()
        )
        return float(cleaned)
    except:
        return None


def categorize_pricing_text(pricing_text: str):
    """
    Maak van pricing-tekst een eenvoudige prijsklasse.
    Retourneert (code, label):
    - code → naar Metric.value
    - label → naar Metric.description
    """
    if not pricing_text:
        return 0, "Onbekend"

    text = pricing_text.lower()

    # Gratis / freemium
    if "free" in text or "gratis" in text:
        return 1, "Gratis / Freemium"

    # Simpele nummerdetectie (eerste bedrag in de tekst)
    nums = re.findall(r"(\d+(?:[.,]\d+)?)", text)
    if nums:
        try:
            val = float(nums[0].replace(",", "."))
            if val < 30:
                return 2, "Lage prijsklasse"
            elif val < 100:
                return 3, "Midden segment"
            else:
                return 4, "Hoge prijsklasse"
        except ValueError:
            pass

    # Keywords als fallback
    if "enterprise" in text:
        return 5, "Enterprise"
    if "pro" in text or "business" in text:
        return 3, "Midden / Pro"

    return 0, "Onbekend"


def features_from_company(company):
    """
    Retourneert (aantal_features, 'feature1, feature2, ...').
    Gebaseerd op company.key_features uit het baseline report.
    """
    if company.key_features:
        features = [str(f) for f in company.key_features]
        return len(features), ", ".join(features)
    return 0, "Geen features"


def extract_positive_reviews(company):
    from app.google_reviews import get_google_reviews

    """
    Eerst proberen we officiële Google Reviews API.
    Als er niks is → fallback op AI / tekstdetectie.
    """
    # 1) Google API proberen
    count, label = get_google_reviews(company.name)
    if count > 0:
        return count, label

    # 2) fallback
    text = ((company.traction_signals or "") + " " + (company.ai_summary or "")).lower()
    matches = re.findall(r"(\d+)\s*\+?\s*(?:reviews|review)", text)
    if matches:
        try:
            n = int(matches[0])
            return n, f"{n} positieve reviews (geschat)"
        except:
            pass

    return 0, "Geen reviews gevonden"


def format_funding_for_metric(company):
    """
    Funding voor de Metric-tabel.

    - Als er een numeric funding in de DB staat → gebruik die en toon '€X'.
    - Als er geen numeric is → gebruik tekst uit funding_history / AI als label.
    - In alle gevallen krijgt Metric.value een simpel numeriek code (0 of bedrag).
    """
    # 1) Eerst proberen numeric funding
    if company.funding is not None:
        try:
            val = float(company.funding)
            label = f"€{int(val):,}".replace(",", ".")
            return val, label
        except Exception:
            pass  # val lukt niet → we vallen door naar tekst

    # 2) Geen numeric funding → probeer tekst te vinden
    txt = (company.funding_history or "").strip()
    if not txt:
        # eventueel nog extra context gebruiken
        txt = ((company.traction_signals or "") + " " + (company.ai_summary or "")).strip()

    if not txt:
        txt = "Onbekend"

    # We geven hier 0.0 als numeric code, label = de tekst
    return 0.0, txt


def estimate_hiring_activity(company):
    """
    Eenvoudige indicatie van hiring-activiteit.
    Geeft (code, label) terug.
    """
    text = ((company.traction_signals or "") +
            " " + (company.product_description or "") +
            " " + (company.ai_summary or "")).lower()

    hiring_keywords = [
        "we are hiring", "we're hiring", "join our team",
        "vacatures", "open positions", "careers"
    ]
    if any(k in text for k in hiring_keywords):
        return 3, "Actief aanwervend"

    if company.team_size:
        ts = company.team_size
        if ts >= 200:
            return 2, "Groot team (stabiele hiring)"
        elif ts >= 30:
            return 2, "Groeiend team"
        elif ts < 30:
            return 1, "Klein team"

    return 0, "Onbekend"


def get_or_create_metric(company_id: int, name: str):
    """
    Haal de metric op of maak ze als ze nog niet bestaat.
    """
    metric = Metric.query.filter_by(company_id=company_id, name=name).first()
    if not metric:
        metric = Metric(
            company_id=company_id,
            name=name,
            active=True,
            tracking_frequency="on_change"
        )
        db.session.add(metric)
    return metric

def backfill_historical_metrics(company_id: int, historical_list: list):
    """
    Schrijft door AI gereconstrueerde historiek weg als 'inferred'.
    Verwacht items met:
      name: "TeamSize" | "Funding" | "Pricing" | "Reviews"
      date: "YYYY-MM-DD"
      value: numeriek
      source: "explicit" | "inferred" (we slaan het als 'inferred' op)
    """
    for item in historical_list or []:
        name = item.get("name")
        date_str = item.get("date")
        value = item.get("value")

        if not (name and date_str):
            continue

        try:
            recorded_at = datetime.fromisoformat(date_str)
        except ValueError:
            continue

        num_value = None
        if value is not None:
            try:
                num_value = Decimal(str(value))
            except Exception:
                pass

        hist = MetricHistory(
            company_id=company_id,
            name=name,
            value=num_value,
            recorded_at=recorded_at,
            source="inferred"
        )
        db.session.add(hist)


def track_metric_history(company_id: int, name: str, value, source: str = "snapshot"):
    """
    Slaat een datapunt op voor historiek / grafieken.
    source:
      - "snapshot": live scrape nu
      - "inferred": AI-reconstructie van verleden
    """
    num_value = None
    if value is not None:
        try:
            if isinstance(value, Decimal):
                num_value = value
            else:
                num_value = Decimal(str(value))
        except Exception:
            num_value = None

    hist = MetricHistory(
        company_id=company_id,
        name=name,
        value=num_value,
        source=source
    )
    db.session.add(hist)



def update_company_metrics(company):
    """
    Vul/werk de 5 kernmetrics bij in de Metric-tabel
    op basis van de huidige Company-waarden.
    """

    # 1) Pricing
    price_code, price_label = categorize_pricing_text(company.pricing or "")
    m_pricing = get_or_create_metric(company.company_id, "Pricing")
    m_pricing.value = price_code
    m_pricing.description = price_label
    m_pricing.last_updated = datetime.utcnow()
    track_metric_history(company.company_id, "Pricing", price_code)

    # 2) Features
    feat_count, feat_label = features_from_company(company)
    m_feat = get_or_create_metric(company.company_id, "Features")
    m_feat.value = feat_count
    m_feat.description = feat_label
    m_feat.last_updated = datetime.utcnow()
    track_metric_history(company.company_id, "Features", feat_count)

    # 3) Reviews
    rev_count, rev_label = extract_positive_reviews(company)
    m_rev = get_or_create_metric(company.company_id, "Reviews")
    m_rev.value = rev_count
    m_rev.description = rev_label
    m_rev.last_updated = datetime.utcnow()
    track_metric_history(company.company_id, "Reviews", rev_count)

    # 4) Funding
    fund_val, fund_label = format_funding_for_metric(company)
    m_fund = get_or_create_metric(company.company_id, "Funding")
    m_fund.value = fund_val
    m_fund.description = fund_label
    m_fund.last_updated = datetime.utcnow()
    track_metric_history(company.company_id, "Funding", fund_val)

    # 5) Hiring
    hiring_code, hiring_label = estimate_hiring_activity(company)
    m_hiring = get_or_create_metric(company.company_id, "Hiring")
    m_hiring.value = hiring_code
    m_hiring.description = hiring_label
    m_hiring.last_updated = datetime.utcnow()
    track_metric_history(company.company_id, "Hiring", hiring_code)


# =====================================================
# INDEX
# =====================================================

@bp.route('/')
def index():
    return render_template('index.html')


# =====================================================
# REGISTER
# =====================================================

@bp.route('/register', methods=['GET', 'POST'])
def register():
    message = ""

    if request.method == 'POST':
        username = request.form.get('username')
        email = request.form.get('email')
        password = request.form.get('password')

        if AppUser.query.filter_by(email=email).first():
            return render_template('register.html', message="Dit e-mailadres is al in gebruik.")

        if AppUser.query.filter_by(username=username).first():
            return render_template('register.html', message="Gebruikersnaam bestaat al.")

        new_user = AppUser(
            username=username,
            email=email,
            password_hash=generate_password_hash(password)
        )

        db.session.add(new_user)
        db.session.commit()

        return redirect(url_for('main.login'))

    return render_template('register.html', message=message)

# =====================================================
# INTERNE FUNCTIE: REFRESH ALL COMPANIES (voor scheduler)
# =====================================================

# LET OP: Geen @bp.route decorateur! Deze functie wordt nu direct door APScheduler aangeroepen.
def refresh_all_companies():
    """
    Wordt wekelijks uitgevoerd door APScheduler. 
    Loopt door ALLE bedrijven om hun data, metrics en change events bij te werken.
    """
    from app import db, create_app # Importeer de app context
    
    # Zorg dat de databasebewerkingen binnen de applicatiecontext vallen
    app = create_app()
    with app.app_context():
        
        companies_to_refresh = Company.query.all()
        
        if not companies_to_refresh:
            print("Scheduler: Geen bedrijven gevonden om te verversen.")
            return

        refreshed_count = 0
        
        for company in companies_to_refresh:
            
            if not company.website_url:
                continue

            try:
                # --- SCRAPEN ---
                result = scrape_website(company.website_url)
                
                if result.get("error"):
                    continue

                # Gebruik 'existing' voor duidelijkheid
                existing = company 

                # ============================================
                # 1) STRATEGIC MOVE DETECTION (GEKOPIEERD UIT /SCRAPE)
                # ============================================
                change_events = []

                # ===== FEATURES =====
                old_features = normalize_list(existing.key_features)
                new_features = normalize_list(result.get("key_features"))

                added_features = [f for f in new_features if f not in old_features]
                removed_features = [f for f in old_features if f not in new_features]

                for f in added_features:
                    change_events.append({
                        "event_type": "new_feature",
                        "description": f"Nieuwe feature toegevoegd: {f}"
                    })

                for f in removed_features:
                    change_events.append({
                        "event_type": "removed_feature",
                        "description": f"Feature verwijderd: {f}"
                    })

                # ===== PRICING =====
                old_price = existing.pricing or ""
                new_price = result.get("pricing") or ""

                if not texts_similar(old_price, new_price):
                    if old_price and new_price:
                        change_events.append({
                            "event_type": "pricing_change",
                            "description": f"Pricing gewijzigd van '{old_price}' → '{new_price}'"
                        })
                    elif new_price:
                        change_events.append({
                            "event_type": "pricing_added",
                            "description": f"Pricing toegevoegd: {new_price}"
                        })
                    elif old_price:
                        change_events.append({
                            "event_type": "pricing_removed",
                            "description": "Pricing verwijderd"
                        })

                # ===== PRODUCT DESCRIPTION =====
                old_product = existing.product_description or ""
                new_product = result.get("product_description") or ""

                if not texts_similar(old_product, new_product):
                    change_events.append({
                        "event_type": "product_change",
                        "description": "Productbeschrijving gewijzigd (mogelijke nieuwe productlijn)"
                    })

                # ===== TARGET SEGMENT =====
                old_segment = existing.target_segment or ""
                new_segment = result.get("target_segment") or ""

                if not texts_similar(old_segment, new_segment):
                    change_events.append({
                        "event_type": "segment_change",
                        "description": "Target segment gewijzigd"
                    })
                    
                # Sla alle echte wijzigingen op in de database
                for ev in change_events:
                    db.session.add(ChangeEvent(
                        company_id=existing.company_id,
                        event_type=ev["event_type"],
                        description=ev["description"]
                    ))

                # ----------------------------------------
                # 2) UPDATE BEDRIJFSGEGEVENS
                # ----------------------------------------
                existing.name = result.get("title") or existing.name
                existing.headquarters = result.get("headquarters")
                existing.office_locations = result.get("office_locations")
                existing.team_size = safe_int(result.get("team_size"))
                existing.funding = result.get("funding")
                existing.funding_history = result.get("funding_history")
                existing.traction_signals = result.get("traction_signals")
                existing.ai_summary = result.get("ai_summary")
                existing.value_proposition = result.get("value_proposition")
                existing.product_description = result.get("product_description")
                existing.target_segment = result.get("target_segment")
                existing.pricing = result.get("pricing")
                existing.key_features = result.get("key_features")
                existing.competitors = result.get("competitors")
                
                # 3) METRICS UPDATEN & GESCHIEDENIS TRACKEN
                update_company_metrics(existing)
                backfill_historical_metrics(existing.company_id, result.get("historical_metrics", []))

                # 4) AUDIT LOG
                db.session.add(AuditLog(
                    company_id=existing.company_id,
                    source_name="Scheduled Refresh (APScheduler)",
                    source_url=existing.website_url
                ))
                
                refreshed_count += 1
                
            except Exception as e:
                db.session.rollback() 
                print(f"Scheduler Fout: Fout bij verversen van {company.name}: {e}")
                continue

        # Commit alle updates in één keer 
        db.session.commit()
        print(f"Scheduler: Succesvol {refreshed_count} bedrijven ververst.")

    return # Geen return code of jsonify nodig

# =====================================================
# LOGIN
# =====================================================

@bp.route('/login', methods=['GET', 'POST'])
def login():
    message = ""

    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')

        user = AppUser.query.filter_by(email=email).first()

        if not user or not check_password_hash(user.password_hash, password):
            return render_template('login.html', message="Ongeldig e-mailadres of wachtwoord")

        session['user_id'] = user.user_id
        session['username'] = user.username

        return redirect(url_for('main.dashboard'))

    return render_template('login.html', message=message)


@bp.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('main.login'))

# =====================================================
# DASHBOARD
# =====================================================

@bp.route('/dashboard', methods=['GET', 'POST'])
@login_required
def dashboard():
    user = AppUser.query.get(session["user_id"])
    scrape_result = None

    # -------------------------
    # FORM HANDLING
    # -------------------------
    if request.method == 'POST':
        form_type = request.form.get('form_type')

        # --- WATCHLIST CONFIG
        if form_type == 'watchlist_config':
            session['watchlist_companies'] = [
                int(cid) for cid in request.form.getlist('companies')
            ]
            session['watchlist_metrics'] = [
                m for m in request.form.getlist('metrics') if m in METRIC_OPTIONS
            ]

        # --- COMPETITOR CONFIG
        elif form_type == 'competitor_config':
            session['tracked_competitors'] = request.form.getlist('competitors')

        # --- SCRAPE
        elif form_type == 'scrape':
            url = request.form.get('scrape_url')
            return redirect(url_for('main.scrape') + f"?url={url}")

    # -------------------------
    # DATA LADEN
    # -------------------------
    watchlist_ids = session.get('watchlist_companies', [])
    metrics_selected = session.get('watchlist_metrics', [])
    competitors_selected = session.get('tracked_competitors', [])

    companies_watchlist = (
        Company.query.filter(Company.company_id.in_(watchlist_ids)).all()
        if watchlist_ids else []
    )

    all_companies = Company.query.all()

    # ----------------------------------------
    # ALLE DETECTEERDE COMPETITORS UIT DATABASE
    # ----------------------------------------
    all_detected_competitors = set()

    for comp in all_companies:
        if comp.competitors:
            for c in comp.competitors:
                competitor_name = ""

                if isinstance(c, dict):
                    competitor_name = c.get('name') or c.get('company_name') or ""
                    if not competitor_name:
                        print(f"DEBUG: Found a competitor dict without a 'name': {c}")

                elif isinstance(c, str):
                    competitor_name = c

                if competitor_name and competitor_name.strip() != "":
                    all_detected_competitors.add(competitor_name)

    all_detected_competitors = sorted(all_detected_competitors)

    # ----------------------------------------
    # ALERTS = CHANGEEVENTS VAN GEMONITORDE COMPANIES (WATCHLIST)
    # ----------------------------------------
    if watchlist_ids:
        recent_events_raw = (
            ChangeEvent.query
            .filter(ChangeEvent.company_id.in_(watchlist_ids))
            .order_by(ChangeEvent.detected_at.desc())
            .limit(3)
            .all()
        )
    else:
        recent_events_raw = (
            ChangeEvent.query
            .order_by(ChangeEvent.detected_at.desc())
            .limit(3)
            .all()
        )

    # ✅ N+1 FIX: companies in bulk ophalen
    company_ids = {e.company_id for e in recent_events_raw}
    companies_bulk = (
        Company.query.filter(Company.company_id.in_(company_ids)).all()
        if company_ids else []
    )
    company_name_by_id = {c.company_id: c.name for c in companies_bulk}

    alerts = []
    for e in recent_events_raw:
        alerts.append({
            "id": e.event_id,
            "company_id": e.company_id,
            "company": company_name_by_id.get(e.company_id, "Onbekend bedrijf"),
            "type": e.event_type,
            "description": e.description,
            "time": e.detected_at
        })

    return render_template(
        'dashboard.html',
        user=user,
        scrape_result=scrape_result,
        watchlist=[c.name for c in companies_watchlist],
        alerts=alerts,
        metric_options=METRIC_OPTIONS,
        metrics_selected=metrics_selected,
        competitors_selected=competitors_selected,
        all_competitors=all_detected_competitors,
        companies=all_companies,
        more_alerts_count=max(0, ChangeEvent.query.count() - 3)
    )

# =====================================================
# COMPANY DETAIL
# =====================================================

@bp.route('/company/<int:company_id>')
@login_required
def company_detail(company_id):
    company = Company.query.get_or_404(company_id)

    # --------- WIJZIGINGEN / EVENTS ---------
    events = (ChangeEvent.query
              .filter_by(company_id=company_id)
              .order_by(ChangeEvent.detected_at.desc())
              .all())

    # --------- HISTORIEK VOOR GRAFIEKEN ---------
    def history_series(metric_name: str):
        rows = (MetricHistory.query
                .filter_by(company_id=company_id, name=metric_name)
                .order_by(MetricHistory.recorded_at.asc())
                .all())
        labels = [r.recorded_at.strftime("%Y-%m-%d %H:%M") for r in rows]
        values = [float(r.value) if r.value is not None else None for r in rows]
        return labels, values

    pricing_labels, pricing_values = history_series("Pricing")
    # Hiring chart gebruikt in de template de variabelen hiring_labels/hiring_values,
    # maar we willen daar eigenlijk Team Size tonen.
    hiring_labels, hiring_values = history_series("TeamSize")

    # Fallback: als TeamSize niet in je DB zit, pak dan Hiring
    if not hiring_labels and not hiring_values:
        hiring_labels, hiring_values = history_series("Hiring")

    review_labels, review_values = history_series("Reviews")  # blijft bestaan voor historiek

    # --------- REVIEW DISTRIBUTIE (zoals gisteren) ---------
    def review_distribution(company):
        """
        Geeft altijd exact 5 waarden terug:
        [1★, 2★, 3★, 4★, 5★]
        Simulatie op basis van Google review-count
        """
        dist = [0, 0, 0, 0, 0]

        try:
            from app.google_reviews import get_google_reviews
            count, _label = get_google_reviews(company.name)

            if count and count > 0:
                dist[4] = int(count * 0.55)  # 5★
                dist[3] = int(count * 0.25)  # 4★
                dist[2] = int(count * 0.10)  # 3★
                dist[1] = int(count * 0.05)  # 2★
                dist[0] = max(0, int(count) - sum(dist))  # 1★ rest
        except Exception:
            pass

        return dist

    review_distribution_values = review_distribution(company)

    # --------- PRICING TIER (badge + chart) ---------
    pricing_tier_code = 0
    if pricing_values:
        for v in reversed(pricing_values):
            if v is not None:
                try:
                    pricing_tier_code = int(round(float(v)))
                except Exception:
                    pricing_tier_code = 0
                break

    pricing_tier_code = max(0, min(5, pricing_tier_code))

    tier_labels = {
        0: "Onbekend",
        1: "Gratis / Freemium",
        2: "Lage prijsklasse",
        3: "Midden segment",
        4: "Hoge prijsklasse",
        5: "Enterprise",
    }
    pricing_tier_label = tier_labels.get(pricing_tier_code, "Onbekend")

    # --------- SIMILAR COMPANIES (SECTOR-FILTER) ---------
    from app.similarity import top_similar_companies_in_same_sector
    all_companies = Company.query.all()
    similar = top_similar_companies_in_same_sector(company, all_companies, top_n=5)

    # --------- RENDER ---------
    return render_template(
        "company_detail.html",
        company=company,
        events=events,

        # grafieken
        pricing_labels=pricing_labels,
        pricing_values=pricing_values,
        hiring_labels=hiring_labels,
        hiring_values=hiring_values,

        # reviews
        review_labels=review_labels,
        review_values=review_values,
        review_distribution_values=review_distribution_values,

        # pricing badge
        pricing_tier_label=pricing_tier_label,
        pricing_tier_code=pricing_tier_code,

        # vergelijkbare bedrijven
        similar=similar,
    )



# =====================================================
# WATCHLIST
# =====================================================

@bp.route('/watchlist')
@login_required
def watchlist():
    ids = session.get('watchlist_companies', [])
    try:
        ids = [int(cid) for cid in ids]
    except Exception:
        ids = []

    metrics_selected = session.get('watchlist_metrics', [])
    companies = Company.query.filter(Company.company_id.in_(ids)).all() if ids else []

    comparison_rows = []
    for c in companies:
        metric_values = {}

        for m_label in metrics_selected:
            label_lower = m_label.lower()

            # Haal eventueel een bijhorende Metric op
            metric = (
                Metric.query
                .filter(
                    Metric.company_id == c.company_id,
                    db.func.lower(Metric.name) == label_lower
                )
                .order_by(Metric.last_updated.desc())
                .first()
            )

            # -----------------------------
            # SPECIALE CASE: PRICING
            # -----------------------------
            if label_lower == "pricing":
                company_pricing = (c.pricing or "").strip()

                # Categorie uit Metric (bijv. "Lage prijsklasse")
                metric_label = None
                if metric:
                    if metric.description and metric.description.strip().lower() != "onbekend":
                        metric_label = metric.description.strip()
                    # → FIX: toon '0' NIET
                    elif metric.value is not None and metric.value != 0:
                        metric_label = str(metric.value)

                # Combineer categorie + tekst, MAAR ‘0 – ...’ mag niet meer voorkomen
                if company_pricing and metric_label:
                    display = f"{metric_label} – {company_pricing}"
                elif company_pricing:
                    display = company_pricing
                elif metric_label:
                    display = metric_label
                else:
                    display = "Onbekend"


            # -----------------------------
            # ANDERE METRICS
            # -----------------------------
            else:
                if metric and (metric.description or metric.value is not None):
                    # Gebruik description als die bestaat, anders numeric value
                    display = metric.description or str(metric.value)
                else:
                    # Fallback op basis van company-velden
                    if label_lower == "features":
                        _, display = features_from_company(c)
                    elif label_lower == "reviews":
                        _, display = extract_positive_reviews(c)
                    elif label_lower == "funding":
                        _, display = format_funding_for_metric(c)
                    elif label_lower == "hiring":
                        _, display = estimate_hiring_activity(c)
                    else:
                        display = "–"

            metric_values[m_label] = display

        comparison_rows.append({'company': c, 'metrics': metric_values})

    # Audit logs per bedrijf
    logs_by_company = {
        c.company_id: AuditLog.query.filter_by(company_id=c.company_id)
        .order_by(AuditLog.retrieved_at.desc()).all()
        for c in companies
    }

    return render_template(
        'watchlist.html',
        metrics=metrics_selected,
        rows=comparison_rows,
        companies=companies,
        logs_by_company=logs_by_company
    )


# =====================================================
# EXPORT WATCHLIST AUDIT (CSV / JSON)
# =====================================================

@bp.route('/export-watchlist-audit')
@login_required
def export_watchlist_audit():
    fmt = request.args.get("format", "csv").lower()

    ids = session.get('watchlist_companies', [])
    if not ids:
        return "Geen bedrijven in watchlist.", 400

    try:
        ids = [int(cid) for cid in ids]
    except Exception:
        return "Ongeldige watchlist IDs.", 400

    logs = (
        AuditLog.query
        .filter(AuditLog.company_id.in_(ids))
        .order_by(AuditLog.retrieved_at.desc())
        .all()
    )

    # ✅ N+1 FIX: company names in bulk ophalen
    company_ids = {log.company_id for log in logs}
    companies_bulk = Company.query.filter(Company.company_id.in_(company_ids)).all() if company_ids else []
    company_name_by_id = {c.company_id: c.name for c in companies_bulk}

    # JSON export
    if fmt == "json":
        out = []
        for log in logs:
            out.append({
                "company": company_name_by_id.get(log.company_id, "Onbekend"),
                "source_name": log.source_name,
                "source_url": log.source_url,
                "retrieved_at": log.retrieved_at.isoformat() if log.retrieved_at else None
            })
        return jsonify(out)

    # CSV export
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["Company", "Source Name", "Source URL", "Retrieved At"])

    for log in logs:
        writer.writerow([
            company_name_by_id.get(log.company_id, "Onbekend"),
            log.source_name,
            log.source_url,
            log.retrieved_at.isoformat() if log.retrieved_at else ""
        ])

    return Response(
        output.getvalue(),
        mimetype="text/csv",
        headers={"Content-Disposition": "attachment; filename=audit_export.csv"}
    )



# =====================================================
# COMPANIES OVERVIEW
# =====================================================

@bp.route('/companies', methods=['GET', 'POST'])
@login_required
def companies():
    message = ""

    # Alleen nog watchlist-actie toelaten (geen manual add meer)
    if request.method == 'POST':
        form_type = request.form.get('form_type')

        if form_type == 'add_to_watchlist':
            try:
                cid = int(request.form.get('company_id'))
                wl = session.get('watchlist_companies', [])
                if cid not in wl:
                    wl.append(cid)
                session['watchlist_companies'] = wl
                if not session.get('watchlist_metrics'):
                    session['watchlist_metrics'] = METRIC_OPTIONS
                message = "✔ Toegevoegd aan watchlist"
            except Exception:
                message = "❌ Kon niet aan watchlist toevoegen."

    # --- Sectorfilter uit querystring ---
    selected_sector_id = request.args.get('sector_id', type=int)

    query = Company.query

    if selected_sector_id:
        query = query.filter(Company.sector_id == selected_sector_id)

    companies = query.all()

    # Alle sectoren voor de dropdown
    sectors = Sector.query.order_by(Sector.name.asc()).all()

    return render_template(
        'companies.html',
        companies=companies,
        sectors=sectors,
        selected_sector_id=selected_sector_id,
        message=message
    )


# =====================================================
# EXPORT: One-click company profile (VC analyst)
# =====================================================
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4
from pptx import Presentation
from pptx.util import Inches, Pt

@bp.route('/company/<int:company_id>/export-pdf')
@login_required
def export_pdf(company_id):
    company = Company.query.get_or_404(company_id)

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    # Title
    story.append(Paragraph(f"<b>{company.name}</b>", styles['Title']))
    story.append(Spacer(1, 16))

    fields = [
        ("Website", company.website_url),
        ("Value Proposition", company.value_proposition),
        ("Product Description", company.product_description),
        ("Target Segment", company.target_segment),
        ("Pricing", company.pricing),
        ("Funding", str(company.funding)),
        ("Team Size", str(company.team_size)),
        ("Traction signals", company.traction_signals),
    ]

    for title, value in fields:
        story.append(Paragraph(f"<b>{title}</b><br/>{value or '—'}", styles['BodyText']))
        story.append(Spacer(1, 12))

    # Features
    if company.key_features:
        story.append(Paragraph("<b>Key Features</b>", styles['Heading2']))
        for f in company.key_features:
            story.append(Paragraph(f"- {f}", styles['BodyText']))
        story.append(Spacer(1, 12))

    # Competitors
    if company.competitors:
        story.append(Paragraph("<b>Competitors</b>", styles['Heading2']))
        for comp in company.competitors:
            story.append(Paragraph(f"- {comp}", styles['BodyText']))

    doc.build(story)
    buffer.seek(0)

    return Response(
        buffer,
        mimetype="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={company.name}_Report.pdf"}
    )

@bp.route('/company/<int:company_id>/export-slides')
@login_required
def export_slides(company_id):
    company = Company.query.get_or_404(company_id)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # Title Page
    slide.shapes.title.text = company.name
    slide.placeholders[1].text = "Baseline Analysis"

    # Content slides
    def add_slide(title, content_list):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        for c in content_list:
            body.add_paragraph().text = c

    add_slide("What they do", [
        company.value_proposition or "—",
        company.product_description or "—"
    ])

    add_slide("Target Segment", [company.target_segment or "—"])

    add_slide("Pricing", [company.pricing or "—"])

    if company.key_features:
        add_slide("Key Features", company.key_features)

    if company.competitors:
        add_slide("Competitors", company.competitors)

    # Save file to memory
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return Response(
        buffer,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={company.name}_Slides.pptx"}
    )

@bp.route('/company/<int:company_id>/export')
@login_required
def export_company(company_id):
    c = Company.query.get_or_404(company_id)
    format = request.args.get("format", "json").lower()

    profile = {
        "company_name": c.name,
        "website": c.website_url,
        "headquarters": c.headquarters,
        "office_locations": c.office_locations,
        "team_size": c.team_size,
        "funding": c.funding,
        "funding_history": c.funding_history,
        "value_proposition": c.value_proposition,
        "product_description": c.product_description,
        "target_segment": c.target_segment,
        "pricing": c.pricing,
        "key_features": c.key_features,
        "competitors": c.competitors,
        "traction_signals": c.traction_signals,
        "ai_summary": c.ai_summary,
        "created_at": c.created_at
    }

    # JSON export
    if format == "json":
        return jsonify(profile)

    # CSV export (flat)
    if format == "csv":
        output = io.StringIO()
        writer = csv.writer(output)
        for k, v in profile.items():
            writer.writerow([k, v])
        output.seek(0)
        return Response(
            output.getvalue(),
            mimetype="text/csv",
            headers={"Content-Disposition": f"attachment; filename={c.name}_profile.csv"}
        )

    # Plain-text due-diligence memo
    if format == "txt":
        memo = []
        memo.append(f"=== DUE DILIGENCE PROFILE: {c.name} ===\n")
        memo.append(f"Website: {c.website_url}\n")
        memo.append(f"HQ: {c.headquarters}\n")
        memo.append(f"Office locations: {c.office_locations}\n")
        memo.append(f"Team size: {c.team_size}\n")
        memo.append(f"Funding: {c.funding}\n")
        memo.append(f"Funding history: {c.funding_history}\n")
        memo.append("\n--- PRODUCT & MARKET ---\n")
        memo.append(f"Value proposition: {c.value_proposition}\n")
        memo.append(f"Product description: {c.product_description}\n")
        memo.append(f"Target segment: {c.target_segment}\n")
        memo.append(f"Pricing: {c.pricing}\n")
        memo.append("\nKey Features:\n")
        memo.extend([f" - {k}" for k in (c.key_features or [])])
        memo.append("\nCompetitors:\n")
        memo.extend([f" - {k}" for k in (c.competitors or [])])
        memo.append("\nTraction signals:\n")
        memo.append(c.traction_signals or "Geen signalen")
        memo.append("\nAI summary:\n")
        memo.append(c.ai_summary or "")

        return Response(
            "\n".join(memo),
            mimetype="text/plain",
            headers={"Content-Disposition": f"attachment; filename={c.name}_memo.txt"}
        )

    return "Unsupported format", 400


# =====================================================
# SCRAPE PAGE (AI + AUTO SAVE + CHANGE DETECTION)
# =====================================================

@bp.route('/scrape', methods=['GET', 'POST'])
@login_required
def scrape():
    # URL uit querystring (bijv. quick scrape vanuit dashboard)
    url_from_query = request.args.get("url")

    # alle sectoren ophalen voor dropdown
    sectors = Sector.query.order_by(Sector.name.asc()).all()

    # GET zonder url → lege pagina met formulier
    if request.method == 'GET' and not url_from_query:
        return render_template('scrape.html', result=None, sectors=sectors)

    # URL kan uit querystring (GET) of uit formulier (POST) komen
    url = url_from_query or request.form.get('url')

    # DEBUG: toon request.form bij POST
    if request.method == 'POST':
        print("DEBUG scrape POST form:", dict(request.form))

    if not url:
        return render_template('scrape.html', result={'error': "Geen URL opgegeven."}, sectors=sectors)

    # sector_id komt alleen uit formulier (POST)
    sector_id_raw = request.form.get('sector_id')
    try:
        sector_id = int(sector_id_raw) if sector_id_raw else None
    except ValueError:
        sector_id = None

    print("DEBUG gekozen sector_id:", sector_id)

    # --- SCRAPEN ---
    result = scrape_website(url)
    if result.get("error"):
        return render_template('scrape.html', result=result, sectors=sectors)

    # --- CHECK OF BEDRIJF BESTAAT ---
    existing = Company.query.filter_by(website_url=url).first()

    # ============================================
    # UPDATE BESTAAND BEDRIJF
    # ============================================
    if existing:
        print("DEBUG: bestaand bedrijf, id =", existing.company_id)

        # STRATEGIC MOVE DETECTION (ongewijzigd)
        change_events = []

        old_features = normalize_list(existing.key_features)
        new_features = normalize_list(result.get("key_features"))
        added_features = [f for f in new_features if f not in old_features]
        removed_features = [f for f in old_features if f not in new_features]

        for f in added_features:
            change_events.append({
                "event_type": "new_feature",
                "description": f"Nieuwe feature toegevoegd: {f}"
            })

        for f in removed_features:
            change_events.append({
                "event_type": "removed_feature",
                "description": f"Feature verwijderd: {f}"
            })

        old_price = existing.pricing or ""
        new_price = result.get("pricing") or ""
        if not texts_similar(old_price, new_price):
            if old_price and new_price:
                change_events.append({
                    "event_type": "pricing_change",
                    "description": f"Pricing gewijzigd van '{old_price}' → '{new_price}'"
                })
            elif new_price:
                change_events.append({
                    "event_type": "pricing_added",
                    "description": f"Pricing toegevoegd: {new_price}"
                })
            elif old_price:
                change_events.append({
                    "event_type": "pricing_removed",
                    "description": "Pricing verwijderd"
                })

        old_product = existing.product_description or ""
        new_product = result.get("product_description") or ""
        if not texts_similar(old_product, new_product):
            change_events.append({
                "event_type": "product_change",
                "description": "Productbeschrijving gewijzigd (mogelijke nieuwe productlijn)"
            })

        old_segment = existing.target_segment or ""
        new_segment = result.get("target_segment") or ""
        if not texts_similar(old_segment, new_segment):
            change_events.append({
                "event_type": "segment_change",
                "description": "Target segment gewijzigd"
            })

        for ev in change_events:
            db.session.add(ChangeEvent(
                company_id=existing.company_id,
                event_type=ev["event_type"],
                description=ev["description"]
            ))

        # --- BEDRIJFSGEGEVENS UPDATEN ---
        existing.name = result.get("title") or existing.name
        existing.headquarters = result.get("headquarters")
        existing.office_locations = result.get("office_locations")
        existing.team_size = safe_int(result.get("team_size"))
        existing.funding = result.get("funding")
        existing.funding_history = result.get("funding_history")
        existing.traction_signals = result.get("traction_signals")
        existing.ai_summary = result.get("ai_summary")
        existing.value_proposition = result.get("value_proposition")
        existing.product_description = result.get("product_description")
        existing.target_segment = result.get("target_segment")
        existing.pricing = result.get("pricing")
        existing.key_features = result.get("key_features")
        existing.competitors = result.get("competitors")

        # 🔥 HIER sector_id ZETTEN
        if sector_id is not None:
            print("DEBUG: existing.company.sector_id vóór:", existing.sector_id)
            existing.sector_id = sector_id
            print("DEBUG: existing.company.sector_id ná:", existing.sector_id)

        # METRICS + HISTORIEK
        update_company_metrics(existing)
        historical = result.get("historical_metrics", [])
        backfill_historical_metrics(existing.company_id, historical)

        db.session.commit()
        print("DEBUG: commit gedaan (bestaand bedrijf)")

        return redirect(url_for('main.company_detail', company_id=existing.company_id))

    # ============================================
    # NIEUW BEDRIJF
    # ============================================
    new_company = Company(
        name=result.get("title") or "Onbekend bedrijf",
        website_url=url,
        headquarters=result.get("headquarters"),
        office_locations=result.get("office_locations"),
        team_size=safe_int(result.get("team_size")),
        funding=result.get("funding"),
        funding_history=result.get("funding_history"),
        traction_signals=result.get("traction_signals"),
        ai_summary=result.get("ai_summary"),
        value_proposition=result.get("value_proposition"),
        product_description=result.get("product_description"),
        target_segment=result.get("target_segment"),
        pricing=result.get("pricing"),
        key_features=result.get("key_features"),
        competitors=result.get("competitors"),
    )

    if sector_id is not None:
        print("DEBUG: new_company.sector_id vóór:", new_company.sector_id)
        new_company.sector_id = sector_id
        print("DEBUG: new_company.sector_id ná:", new_company.sector_id)

    db.session.add(new_company)
    db.session.flush()  # zodat new_company.company_id bestaat

    db.session.add(AuditLog(
        company_id=new_company.company_id,
        source_name="Scraper + AI",
        source_url=url
    ))

    update_company_metrics(new_company)
    historical = result.get("historical_metrics", [])
    backfill_historical_metrics(new_company.company_id, historical)

    db.session.commit()
    print("DEBUG: commit gedaan (nieuw bedrijf)")

    return redirect(url_for('main.company_detail', company_id=new_company.company_id))




# =====================================================
# GLOBAL AUDIT LOG OVERVIEW (Compliance)
# =====================================================

@bp.route('/audit-logs')
@admin_required
def audit_logs():
    logs = AuditLog.query.order_by(AuditLog.retrieved_at.desc()).all()

    # ✅ N+1 FIX: company names in bulk ophalen
    company_ids = {log.company_id for log in logs}
    companies_bulk = Company.query.filter(Company.company_id.in_(company_ids)).all() if company_ids else []
    company_name_by_id = {c.company_id: c.name for c in companies_bulk}

    enriched_logs = []
    for log in logs:
        enriched_logs.append({
            "company": company_name_by_id.get(log.company_id, "Onbekend"),
            "company_id": log.company_id,
            "source_name": log.source_name,
            "source_url": log.source_url,
            "retrieved_at": log.retrieved_at
        })

    return render_template("audit_logs.html", logs=enriched_logs)

# ====================================================
# EXPORT ALL AUDIT LOGS (CSV / JSON)
# ====================================================

@bp.route('/audit-logs/export')
@admin_required
def export_all_audit_logs():
    fmt = request.args.get("format", "csv").lower()

    logs = AuditLog.query.order_by(AuditLog.retrieved_at.desc()).all()

    # ✅ N+1 FIX: company names in bulk ophalen
    company_ids = {log.company_id for log in logs}
    companies_bulk = Company.query.filter(Company.company_id.in_(company_ids)).all() if company_ids else []
    company_name_by_id = {c.company_id: c.name for c in companies_bulk}

    # JSON export
    if fmt == "json":
        out = []
        for log in logs:
            out.append({
                "company": company_name_by_id.get(log.company_id, "Onbekend"),
                "source_name": log.source_name,
                "source_url": log.source_url,
                "retrieved_at": log.retrieved_at.isoformat() if log.retrieved_at else None
            })
        return jsonify(out)

    # CSV export
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["Company", "Source Name", "Source URL", "Retrieved At"])

    for log in logs:
        writer.writerow([
            company_name_by_id.get(log.company_id, "Onbekend"),
            log.source_name,
            log.source_url,
            log.retrieved_at.isoformat() if log.retrieved_at else ""
        ])

    return Response(
        output.getvalue(),
        mimetype="text/csv",
        headers={"Content-Disposition": "attachment; filename=all_audit_logs.csv"}
    )

# =====================================================
# WEEKLY MAIL SETTINGS
# =====================================================

@bp.route("/weekly-mail-settings")
@login_required
def weekly_mail_settings():
    user = AppUser.query.get(session["user_id"])
    return render_template("weekly_mail.html", user=user)


@bp.route("/update-weekly-mail", methods=["POST"])
@login_required
def update_weekly_mail():
    user = AppUser.query.get(session["user_id"])
    user.weekly_digest = request.form.get("digest") == "on"
    db.session.commit()

    return redirect(url_for("main.weekly_mail_settings"))


# =====================================================
# ALL ALERTS PAGE
# =====================================================

@bp.route("/all-alerts")
@login_required
def all_alerts():
    events = ChangeEvent.query.order_by(
        ChangeEvent.detected_at.desc()
    ).all()

    # ✅ N+1 FIX: company names in bulk ophalen
    company_ids = {e.company_id for e in events}
    companies_bulk = Company.query.filter(Company.company_id.in_(company_ids)).all() if company_ids else []
    company_name_by_id = {c.company_id: c.name for c in companies_bulk}

    alerts = []
    for e in events:
        alerts.append({
            "company_id": e.company_id,
            "company": company_name_by_id.get(e.company_id, "Onbekend bedrijf"),
            "type": e.event_type,
            "description": e.description,
            "time": e.detected_at
        })

    return render_template(
        "all_alerts.html",
        alerts=alerts,
        company=None
    )

# =====================================================
# API: FEED VAN ALLE CHANGE EVENTS (voor BI / analytics)
# =====================================================

@bp.route("/api/events")
@admin_required
def api_events():
    """
    JSON feed van alle gedetecteerde events:
    - pricing changes
    - feature changes
    - (later) funding, hiring, segment changes, geo-expansion...
    Optionele filters:
      ?company_id=...
      ?type=pricing_change
      ?since=2025-12-01
    """

    # ✅ N+1 FIX: query levert meteen Company.name mee
    query = (
        db.session.query(ChangeEvent, Company.name)
        .join(Company, ChangeEvent.company_id == Company.company_id)
    )

    company_id = request.args.get("company_id", type=int)
    if company_id:
        query = query.filter(ChangeEvent.company_id == company_id)

    event_type = request.args.get("type")
    if event_type:
        query = query.filter(ChangeEvent.event_type == event_type)

    since_str = request.args.get("since")
    if since_str:
        try:
            since_date = datetime.fromisoformat(since_str)
            query = query.filter(ChangeEvent.detected_at >= since_date)
        except ValueError:
            pass

    rows = query.order_by(ChangeEvent.detected_at.desc()).all()

    out = []
    for e, company_name in rows:
        out.append({
            "event_id": e.event_id,
            "company_id": e.company_id,
            "company_name": company_name or "Onbekend",
            "event_type": e.event_type,
            "description": e.description,
            "detected_at": e.detected_at.isoformat() if e.detected_at else None
        })

    return jsonify(out)

# =====================================================
# DELETE COMPANY 
# =====================================================

@bp.route("/company/<int:company_id>/delete", methods=["POST"])
@login_required
def delete_company(company_id):
    company = Company.query.get(company_id)
    if not company:
        return "Bedrijf niet gevonden", 404

    # Verwijder het bedrijf en commit de transactie
    db.session.delete(company)
    db.session.commit()

    # Optioneel: In een productie-omgeving zou u ook Audit Logs,
    # Metrics, Change Events, en Metric History gerelateerd aan
    # dit bedrijf moeten verwijderen om weesgegevens te voorkomen.

    # Terugsturen naar het bedrijvenoverzicht
    return redirect(url_for("main.companies"))

# =====================================================
# COMPANY ALERTS 
# =====================================================

@bp.route('/company/<int:company_id>/alerts')
@login_required
def company_alerts(company_id):
    company = Company.query.get_or_404(company_id)
    
    # Haal alle events op voor dit specifieke bedrijf
    events = ChangeEvent.query.filter_by(company_id=company_id).order_by(
        ChangeEvent.detected_at.desc()
    ).all()

    alerts = []
    for e in events:
        alerts.append({
            "company_id": e.company_id,
            "company": company.name,
            "type": e.event_type,
            "description": e.description,
            "time": e.detected_at
        })
    
    # Gebruik dezelfde template als /all-alerts, maar geef het bedrijf mee
    return render_template(
        "all_alerts.html",
        alerts=alerts,
        company=company # Nu weet de template welk bedrijf het betreft
    )